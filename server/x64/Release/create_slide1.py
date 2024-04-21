from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import datetime
import os
import psycopg2

# Define color constants
RED = RGBColor(211, 0, 0)
GREEN = RGBColor(0, 211, 0)
LIGHTGREY = RGBColor(211, 211, 211)
BLACK = RGBColor(0, 0, 0)

# Database connection parameters
db_host = 'localhost'
db_port = 5432
db_name = 'calendar'
db_user = 'postgres'
db_password = 'password'

# Connect to the database
conn = psycopg2.connect(host=db_host, port=db_port, database=db_name, user=db_user, password=db_password)
cursor = conn.cursor()

# Get current date and time
current_datetime = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
current_date = datetime.datetime.now().strftime('%Y-%m-%d')
current_time = datetime.datetime.now().replace(minute=30 if datetime.datetime.now().minute >= 30 else 0, second=0, microsecond=0).strftime('%H:%M:%S')

# Create a new PowerPoint presentation
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Study Room Information"
subtitle.text = current_datetime

# Add study aids slide
slide_layout = prs.slide_layouts[5]  # Title and content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Study Aids"

# Fetch distinct room options from the database where option_unavailable is false
cursor.execute("SELECT DISTINCT option_name FROM room_options WHERE option_unavailable = 'false'")
room_options = [option[0] for option in cursor.fetchall()]

# Create a grid for study aids options
num_options = len(room_options)
num_cols = 4
num_rows = (num_options // num_cols) + (1 if num_options % num_cols != 0 else 0)
cell_width = Inches(1.5)
cell_height = Inches(0.75)
outline_width = Inches(0.05)

# Calculate the width and height of the grid
grid_width = num_cols * cell_width
grid_height = num_rows * cell_height

# Calculate the left and top margins to center the grid
left_margin = (Inches(10) - grid_width) / 2
top_margin = (Inches(7.5) - grid_height) / 2

# Add cells to the slide
for i, option_name in enumerate(room_options):
    row = i // num_cols
    col = i % num_cols
    left = left_margin + col * cell_width
    top = top_margin + row * cell_height
    shape = slide.shapes.add_shape(
        1, left, top, cell_width, cell_height,    # shape type=rectangle, left, top, width, height
    )
    text_frame = shape.text_frame
    p = text_frame.add_paragraph()
    p.text = option_name
    p.alignment = 1  # Center alignment
    
    # Set the font color based on key availability
    font_color = BLACK
    for run in p.runs:
        run.font.color.rgb = font_color

    # Set cell properties
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHTGREY
    shape.line.width = outline_width
    shape.line.color.rgb = BLACK
    text_frame.margin_top = Inches(0.1)

# Create a new slide for the grid
slide_layout = prs.slide_layouts[5]  # Title and content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Room Availability"

# Fetch all rooms from the room_keys table, sorted alphabetically
cursor.execute("SELECT DISTINCT room FROM room_keys ORDER BY room")
rooms = [room[0] for room in cursor.fetchall()]

# Create a grid with room availability information
num_rooms = len(rooms)
num_cols = 4
num_rows = (num_rooms // num_cols) + (1 if num_rooms % num_cols != 0 else 0)
cell_width = Inches(1.5)
cell_height = Inches(0.75)
outline_width = Inches(0.05)

# Calculate the width and height of the grid
grid_width = num_cols * cell_width
grid_height = num_rows * cell_height

# Calculate the left and top margins to center the grid
left_margin = (Inches(10) - grid_width) / 2
top_margin = (Inches(7.5) - grid_height) / 2

# Add cells to the slide
for i, room in enumerate(rooms):
    row = i // num_cols
    col = i % num_cols
    left = left_margin + col * cell_width
    top = top_margin + row * cell_height
    shape = slide.shapes.add_shape(
        1, left, top, cell_width, cell_height,    # shape type=rectangle, left, top, width, height
    )
    text_frame = shape.text_frame
    p = text_frame.add_paragraph()
    p.text = room
    p.alignment = 1  # Center alignment

    # Check if the room key is available and not blocked for the current date and timeslot
    cursor.execute("SELECT key_available FROM room_keys WHERE room = %s AND date = %s AND timeslot = %s", (room, current_date, current_time))
    key_result = cursor.fetchone()

    cursor.execute("SELECT is_blocked FROM availabilities WHERE room = %s AND date = %s AND timeslot = %s", (room, current_date, current_time))
    block_result = cursor.fetchone()

    if key_result is not None and block_result is not None:
        key_available = key_result[0]
        is_blocked = block_result[0]
    else:
        key_available = False
        is_blocked = False

    # Set the font color based on key availability and block status
    font_color = GREEN if key_available and not is_blocked else RED
    for run in p.runs:
        run.font.color.rgb = font_color

    # Set cell properties
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHTGREY
    shape.line.width = outline_width
    shape.line.color.rgb = BLACK 
    text_frame.margin_top = Inches(0.1)

# Save the presentation
dir_path = os.path.dirname(os.path.realpath(__file__))
pptx_file_path = os.path.join(dir_path, 'presentation', 'complete_presentation.pptx')
prs.save(pptx_file_path)




# Image arc

# Create a new image for the title slide
title_slide_image = Image.new('RGB', (1920, 1080), color='white')
title_draw = ImageDraw.Draw(title_slide_image)
title_font = ImageFont.truetype("arial", 100)  # Adjust font size as needed
title_text = "Study Room Information\n" + current_datetime

# Draw the text in the center
title_draw.text((400, 400), title_text, fill='black', font=title_font)

# Save the title slide image
title_slide_image.save('slide_1.png')

# Create a new image for the study aids slide
study_aids_image = Image.new('RGB', (1920, 1080), color='white')
study_aids_draw = ImageDraw.Draw(study_aids_image)
study_aids_font_title = ImageFont.truetype("arial", 100)  # Adjust font size and font file as needed
study_aids_font = ImageFont.truetype("arial", 60)  # Adjust font size and font file as needed
study_aids_draw.text((720, 400), "Study Aids", fill='black', font=study_aids_font_title)

# Fetch distinct room options from the database where option_unavailable is false
cursor.execute("SELECT DISTINCT option_name FROM room_options WHERE option_unavailable = 'false'")
room_options = [option[0] for option in cursor.fetchall()]

# Create a grid for study aids options
num_options = len(room_options)
num_cols = 4
num_rows = (num_options // num_cols) + (1 if num_options % num_cols != 0 else 0)
cell_width = 400
cell_height = 100
outline_width = 5
cell_padding = 10  # Padding between text and cell border

# Calculate the width and height of the grid
grid_width = num_cols * cell_width
grid_height = num_rows * cell_height

# Calculate the left and top margins to center the grid with a 100px right shift
left_margin = (1920 - grid_width) / 2
top_margin = (1080 - grid_height) / 2 + 100  # Start below the title

# Add cells to the image
for i, option_name in enumerate(room_options):
    row = i // num_cols
    col = i % num_cols
    left = left_margin + col * cell_width
    top = top_margin + row * cell_height
    # Draw light grey rectangle as background
    study_aids_draw.rectangle([left, top, left + cell_width, top + cell_height], fill='lightgrey', outline='black', width=outline_width)
    study_aids_draw.text((left + cell_padding, top + cell_padding), option_name, fill='black', font=study_aids_font)

# Save the study aids image
study_aids_image.save('slide_2.png')

# Create a new image for the room availability slide
room_availability_image = Image.new('RGB', (1920, 1080), color='white')
room_availability_draw = ImageDraw.Draw(room_availability_image)
room_availability_font_title = ImageFont.truetype("arial", 100)  # Adjust font size and font file as needed
room_availability_font = ImageFont.truetype("arial", 60)  # Adjust font size and font file as needed

# Add title
title_text = "Room Availability"
title_text_width = len(title_text) * 30  # Estimate width based on character count
title_left = (1920 - title_text_width) / 2
room_availability_draw.text((570, 400), title_text, fill='black', font=room_availability_font_title)

# Fetch all rooms from the room_keys table, sorted alphabetically
cursor.execute("SELECT DISTINCT room FROM room_keys ORDER BY room")
rooms = [room[0] for room in cursor.fetchall()]

# Create a grid with room availability information
num_rooms = len(rooms)
num_cols = 4
num_rows = (num_rooms // num_cols) + (1 if num_rooms % num_cols != 0 else 0)
cell_width = 400
cell_height = 100
outline_width = 5

# Calculate the width and height of the grid
grid_width = num_cols * cell_width
grid_height = num_rows * cell_height

# Calculate the left and top margins to center the grid
left_margin = (1920 - grid_width) / 2
top_margin = (1080 - grid_height) / 2 + 100  # Start below the title

# Add cells to the image
for i, room in enumerate(rooms):
    row = i // num_cols
    col = i % num_cols
    left = left_margin + col * cell_width
    top = top_margin + row * cell_height

    # Check if the room key is available and not blocked for the current date and timeslot
    cursor.execute("SELECT key_available FROM room_keys WHERE room = %s AND date = %s AND timeslot = %s", (room, current_date, current_time))
    key_result = cursor.fetchone()

    cursor.execute("SELECT is_blocked FROM availabilities WHERE room = %s AND date = %s AND timeslot = %s", (room, current_date, current_time))
    block_result = cursor.fetchone()

    if key_result is not None and block_result is not None:
        key_available = key_result[0]
        is_blocked = block_result[0]
    else:
        key_available = False
        is_blocked = False

    # Set the fill color and outline color based on key availability and block status
    fill_color = 'green' if key_available and not is_blocked else 'red'
    outline_color = 'black'
    room_availability_draw.rectangle([left, top, left + cell_width, top + cell_height], fill=fill_color, outline=outline_color, width=outline_width)

    # Add text to the cell
    room_availability_draw.text((left + 10, top + 10), room, fill='black', font=room_availability_font)

# Save the room availability image
room_availability_image.save('slide_3.png')

# Close database connection
cursor.close()
conn.close()